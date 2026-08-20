"""MarkItDown PPTX converter with font-size heading promotion.

MarkItDown only promotes shapes that are real title placeholders
(``slide.shapes.title``). Decks designed with plain text boxes instead of
placeholders — very common — come out flat. This converter additionally
promotes text boxes whose font size is at least 1.2x the slide's median
body size to ``## `` headings, mirroring the PDF converter.
"""

from __future__ import annotations

import contextlib
import statistics

import pptx
from markitdown._base_converter import DocumentConverterResult
from markitdown.converters import PptxConverter

_BODY_RATIO = 1.2
_MAX_TITLE_CHARS = 120


def _shape_max_size(shape) -> float:
    """Largest explicit run size in a shape's text frame (0 when unknown)."""
    if not shape.has_text_frame:
        return 0.0
    sizes: list[float] = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            size = run.font.size
            if size is not None:
                sizes.append(float(size.pt))
    return max(sizes) if sizes else 0.0


def _iter_shapes(slide):
    for shape in slide.shapes:
        yield shape
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape)


class HeadingPptxConverter(PptxConverter):
    """PPTX converter that promotes large-font text boxes to headings."""

    def convert(self, file_stream, stream_info, **kwargs):
        presentation = pptx.Presentation(file_stream)
        md_content = ""
        for slide_num, slide in enumerate(presentation.slides, start=1):
            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title

            sizes = {
                shape.shape_id: _shape_max_size(shape) for shape in _iter_shapes(slide)
            }
            known = [size for size in sizes.values() if size > 0]
            body = statistics.median(known) if known else None
            threshold = body * _BODY_RATIO if body else None

            def get_shape_content(
                shape, *, title=title, threshold=threshold, sizes=sizes, **kwargs
            ):
                nonlocal md_content
                # Pictures
                if self._is_picture(shape):
                    alt_text = shape.name
                    with contextlib.suppress(Exception):
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    alt_text = _re_sub_whitespace(alt_text)
                    if kwargs.get("keep_data_uris", False):
                        with contextlib.suppress(Exception):
                            content_type, b64_string = _image_data_uri(shape)
                            md_content += (
                                f"\n![{alt_text}](data:{content_type};base64,{b64_string})\n"
                            )
                            return
                    md_content += f"\n[Image: {alt_text or 'Picture'}]\n"
                    return

                # Tables
                if self._is_table(shape):
                    md_content += self._convert_table_to_markdown(shape.table, **kwargs)
                    return

                # Charts
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)
                    return

                # Text areas
                if shape.has_text_frame:
                    text = shape.text.lstrip()
                    if shape == title:
                        md_content += "# " + text + "\n"
                    elif (
                        threshold is not None
                        and sizes.get(shape.shape_id, 0.0) >= threshold
                        and len(text) <= _MAX_TITLE_CHARS
                    ):
                        md_content += "## " + text + "\n"
                    else:
                        md_content += shape.text + "\n"
                    return

                # Group shapes
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                    sorted_shapes = sorted(
                        shape.shapes,
                        key=lambda x: (
                            float("-inf") if not x.top else x.top,
                            float("-inf") if not x.left else x.left,
                        ),
                    )
                    for subshape in sorted_shapes:
                        get_shape_content(subshape, **kwargs)

            sorted_shapes = sorted(
                slide.shapes,
                key=lambda x: (
                    float("-inf") if not x.top else x.top,
                    float("-inf") if not x.left else x.left,
                ),
            )
            for shape in sorted_shapes:
                get_shape_content(shape, **kwargs)

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(markdown=md_content.strip())


def _image_data_uri(shape):
    """Return (content_type, base64) for a picture shape, raising on failure."""
    import base64

    image = shape.image
    b64_string = base64.b64encode(image.blob).decode("utf-8")
    return image.content_type or "image/png", b64_string


def _re_sub_whitespace(text: str) -> str:
    import re

    text = re.sub(r"[\r\n\[\]]", " ", text)
    return re.sub(r"\s+", " ", text).strip()