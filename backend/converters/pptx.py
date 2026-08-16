"""PPTX converter built on python-pptx.

Iterates slides in order, extracts titles as headings, body text with list
structure, tables, pictures, hyperlinks, chart data and speaker notes, and
inserts slide boundaries between slides.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from converters.base import BaseConverter, CorruptFileError
from document_model.blocks import (
    BulletListBlock,
    HeadingBlock,
    ImageBlock,
    ParagraphBlock,
    QuoteBlock,
    SlideBreakBlock,
    TableBlock,
    TableCell,
    TableRow,
    TextRun,
)
from document_model.document import Document


class PptxConverter(BaseConverter):
    format = "pptx"
    extensions = (".pptx",)

    def convert(self) -> Document:
        path = Path(self.context.source_path)
        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise CorruptFileError("The PPTX file could not be opened or is corrupted.") from exc

        doc_model = self._build_document(self._metadata(prs))
        blocks: list = []
        stats = doc_model.stats
        total = len(prs.slides)

        for index, slide in enumerate(prs.slides, start=1):
            self.context.progress(
                "extract", index, total, f"Extracting slide {index} of {total}"
            )
            blocks.append(SlideBreakBlock(slide_number=index))
            self._convert_slide(slide, blocks, doc_model)

        stats.slides = total
        doc_model.blocks = blocks
        doc_model.stats = stats
        return doc_model

    def _metadata(self, prs: Presentation) -> dict:
        core = prs.core_properties
        return {
            "title": core.title,
            "author": core.author,
            "subject": core.subject,
            "keywords": core.keywords,
            "comments": core.comments,
            "created": core.created,
            "modified": core.modified,
            "last_modified_by": core.last_modified_by,
        }

    def _convert_slide(self, slide, blocks: list, doc_model: Document) -> None:
        stats = doc_model.stats
        title_done = False

        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE) and not title_done:
                    title_done = True
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        runs = self._runs_from_frame(shape.text_frame)
                        blocks.append(HeadingBlock(level=1, content=runs))
                        stats.headings += 1
                    continue

            if shape.shape_type == MSO_SHAPE_TYPE.TABLE and shape.has_table:
                self._convert_table(shape.table, blocks, doc_model)
                continue

            if shape.has_chart:
                self._convert_chart(shape, blocks, doc_model)
                continue

            if shape.shape_type in (
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.LINKED_PICTURE,
                MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
            ):
                self._convert_picture(shape, blocks, doc_model)
                continue

            if shape.has_text_frame and shape.text_frame.text.strip():
                self._convert_text_frame(shape.text_frame, blocks, doc_model)

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                runs = [TextRun(text=line.strip()) for line in notes_text.splitlines() if line.strip()]
                blocks.append(QuoteBlock(content=runs))
                stats.paragraphs += 1

    def _convert_text_frame(self, frame, blocks: list, doc_model: Document) -> None:
        stats = doc_model.stats
        grouped: list[tuple[int, bool, list[TextRun]]] = []
        for paragraph in frame.paragraphs:
            text = "".join(r.text for r in paragraph.runs)
            if not text.strip():
                continue
            level = paragraph.level
            is_bullet = paragraph._pPr is not None and (
                paragraph._pPr.get("buChar") is not None or paragraph._pPr.get("buAutoNum") is not None
            )
            runs = self._runs_from_paragraph(paragraph)
            grouped.append((level, is_bullet, runs))

        current_list: list[tuple[int, bool, list[TextRun]]] | None = None
        for level, is_bullet, runs in grouped:
            if is_bullet and current_list is not None and current_list[0][0] == level:
                current_list[1].append(runs)
                continue
            if is_bullet:
                if current_list is not None:
                    self._flush_list(current_list, blocks, doc_model)
                current_list = (level, [runs])
                continue
            if current_list is not None:
                self._flush_list(current_list, blocks, doc_model)
                current_list = None
            stats.paragraphs += 1
            blocks.append(ParagraphBlock(content=runs))

        if current_list is not None:
            self._flush_list(current_list, blocks, doc_model)

    def _flush_list(self, entry: tuple[int, list], blocks: list, doc_model: Document) -> None:
        level, items = entry
        prefix = "\t" * min(level, 4)
        items = [[TextRun(text=prefix), *item] if prefix else item for item in items]
        last = blocks[-1] if blocks else None
        if isinstance(last, BulletListBlock) and last.metadata.get("level") == level:
            last.items.extend(items)
        else:
            block = BulletListBlock(items=items)
            block.metadata["level"] = level
            blocks.append(block)
            doc_model.stats.lists += 1
        doc_model.stats.paragraphs += len(items)

    def _runs_from_frame(self, frame) -> list[TextRun]:
        runs: list[TextRun] = []
        for paragraph in frame.paragraphs:
            runs.extend(self._runs_from_paragraph(paragraph))
            runs.append(TextRun(text="\n"))
        return runs

    def _runs_from_paragraph(self, paragraph) -> list[TextRun]:
        runs: list[TextRun] = []
        for run in paragraph.runs:
            if not run.text:
                continue
            href = None
            if run.hyperlink and run.hyperlink.address:
                href = run.hyperlink.address
            font = run.font
            runs.append(
                TextRun(
                    text=run.text,
                    bold=bool(font.bold),
                    italic=bool(font.italic),
                    underline=bool(font.underline),
                    href=href,
                )
            )
        return runs

    def _convert_table(self, table, blocks: list, doc_model: Document) -> None:
        stats = doc_model.stats
        rows: list[TableRow] = []
        for row in table.rows:
            cells: list[TableCell] = []
            for cell in row.cells:
                runs: list[TextRun] = []
                for paragraph in cell.text_frame.paragraphs:
                    runs.extend(self._runs_from_paragraph(paragraph))
                cells.append(TableCell(content=runs))
            rows.append(TableRow(cells=cells))
        if rows:
            if self.context.settings.convert_tables:
                blocks.append(TableBlock(rows=rows, has_header=True))
                stats.tables += 1
            else:
                for row in rows:
                    values = ["".join(r.text for r in cell.content) for cell in row.cells]
                    stats.paragraphs += 1
                    blocks.append(ParagraphBlock(content=[TextRun(text=" | ".join(values))]))

    def _convert_chart(self, shape, blocks: list, doc_model: Document) -> None:
        """Represent chart data as a table so numeric content is never lost."""
        stats = doc_model.stats
        try:
            chart = shape.chart
            title = ""
            if getattr(chart, "has_title", False) and chart.chart_title is not None:
                title = chart.chart_title.text_frame.text.strip()

            categories: list[str] = []
            series: list[tuple[str, list[float]]] = []
            for plot in chart.plots:
                try:
                    plot_cats = [str(c) for c in plot.categories]
                except Exception:
                    plot_cats = []
                for s in plot.series:
                    name = s.name or ""
                    try:
                        values = [v for v in s.values if isinstance(v, (int, float))]
                    except Exception:
                        values = []
                    if not categories and plot_cats:
                        categories = plot_cats
                    series.append((name, values))

            if not series:
                return
            rows: list[TableRow] = []
            header: list[TableCell] = [TableCell(content=[TextRun(text="Category")])]
            header.extend(TableCell(content=[TextRun(text=name)]) for name, _ in series)
            rows.append(TableRow(cells=header))
            row_count = max((len(values) for _, values in series), default=0)
            for index in range(row_count):
                category_cell = (
                    [TableCell(content=[TextRun(text=categories[index])])]
                    if index < len(categories)
                    else [TableCell(content=[])]
                )
                cells = list(category_cell)
                for _, values in series:
                    value = values[index] if index < len(values) else ""
                    cells.append(TableCell(content=[TextRun(text=self._format_value(value))]))
                rows.append(TableRow(cells=cells))

            if self.context.settings.convert_tables:
                if title:
                    blocks.append(ParagraphBlock(content=[TextRun(text=f"Chart: {title}")]))
                    stats.paragraphs += 1
                blocks.append(TableBlock(rows=rows, has_header=True))
                stats.tables += 1
                stats.charts += 1
            else:
                stats.paragraphs += 1
                summary = f"Chart{': ' + title if title else ''}: " + ", ".join(
                    f"{name or 'series'} = {', '.join(self._format_value(v) for v in values)}"
                    for name, values in series
                )
                blocks.append(ParagraphBlock(content=[TextRun(text=summary)]))
        except Exception:
            doc_model.warnings.append(
                {
                    "code": "chart_extraction_failed",
                    "message": "A chart could not be parsed and was skipped.",
                    "severity": "warning",
                }
            )

    @staticmethod
    def _format_value(value) -> str:
        if isinstance(value, float) and value.is_integer():
            return str(int(value))
        return str(value)

    def _convert_picture(self, shape, blocks: list, doc_model: Document) -> None:
        if not hasattr(shape, "image"):
            return
        try:
            blob = shape.image.blob
            ext = shape.image.ext
        except Exception:
            return
        alt = shape.name or ""
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, "_element"):
            desc = shape._element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr")
            if desc is not None:
                alt = desc.get("descr") or desc.get("name") or alt
        rel = self.context.save_image(blob, ext, alt=alt)
        if rel:
            doc_model.stats.images += 1
            blocks.append(ImageBlock(path=rel, alt=alt))
