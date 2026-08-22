"""Unit tests for the MarkItDown conversion engine adapter."""

from __future__ import annotations

from pathlib import Path

import pytest
from fixtures.make_fixtures import (
    make_deck_pdf,
    make_docx,
    make_pdf,
    make_pptx,
    make_scanned_pdf,
    make_two_column_pdf,
    make_xlsx,
)

from app.schemas.settings import ConversionSettings
from converters.base import ConversionContext
from converters.markitdown import convert_with_markitdown
from document_model.document import Document


def make_context(source: Path, tmp_path: Path, **overrides) -> ConversionContext:
    output_dir = tmp_path / "out"
    output_dir.mkdir(exist_ok=True)
    return ConversionContext(
        source_path=source,
        settings=ConversionSettings(**overrides),
        output_dir=output_dir,
    )


def assert_no_figures(markdown: str) -> None:
    """These fixtures embed no images, so no figure should be referenced.

    Image extraction is on by default, so this now guards against inventing
    references rather than against extraction itself. Figure extraction has its
    own coverage in :class:`TestImageExtraction`.
    """
    assert markdown
    assert "![" not in markdown
    assert "assets/" not in markdown


class TestMarkitdownEngine:
    def test_pdf_conversion(self, tmp_path: Path):
        source = make_pdf(tmp_path / "report.pdf", pages=3)
        context = make_context(source, tmp_path)
        doc = convert_with_markitdown(context)
        assert isinstance(doc, Document)
        assert doc.format == "pdf"
        assert context.markdown_output
        assert_no_figures(context.markdown_output)
        assert doc.stats.headings >= 0
        assert len(doc.warnings) == 0

    def test_docx_conversion(self, tmp_path: Path):
        source = make_docx(tmp_path / "letter.docx")
        context = make_context(source, tmp_path)
        doc = convert_with_markitdown(context)
        assert doc.format == "docx"
        assert context.markdown_output
        assert_no_figures(context.markdown_output)

    def test_pptx_conversion(self, tmp_path: Path):
        source = make_pptx(tmp_path / "deck.pptx", slides=3)
        context = make_context(source, tmp_path)
        doc = convert_with_markitdown(context)
        assert doc.format == "pptx"
        assert context.markdown_output
        assert_no_figures(context.markdown_output)

    def test_xlsx_conversion(self, tmp_path: Path):
        source = make_xlsx(tmp_path / "book.xlsx", sheets=2)
        context = make_context(source, tmp_path)
        doc = convert_with_markitdown(context)
        assert doc.format == "xlsx"
        assert context.markdown_output
        assert_no_figures(context.markdown_output)

    def test_pptx_pictures_become_placeholders_when_extraction_is_off(
        self, tmp_path: Path
    ):
        from io import BytesIO

        from PIL import Image
        from pptx import Presentation

        source = tmp_path / "pics.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Pictures"
        buffer = BytesIO()
        Image.new("RGB", (16, 16), (200, 30, 30)).save(buffer, format="PNG")
        buffer.seek(0)
        slide.shapes.add_picture(buffer, 0, 0)
        prs.save(source)

        context = make_context(source, tmp_path, extract_images=False)
        convert_with_markitdown(context)
        out = context.markdown_output
        assert out
        assert "![" not in out
        # python-pptx auto-names an added picture shape "image.png" when no
        # real name/description was set -- that filename-shaped default is
        # exactly what _looks_like_placeholder_alt now filters out, so the
        # generic "Picture" fallback is used instead of leaking the filename.
        assert "[Image: Picture]" in out

    def test_pptx_blank_notes_are_omitted(self, tmp_path: Path):
        """PowerPoint creates the notes placeholder on every slide regardless
        of whether the author typed anything -- has_notes_slide alone must not
        be enough to emit a "### Notes:" heading with nothing under it."""
        from pptx import Presentation

        source = tmp_path / "notes.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide with blank notes"
        # Accessing notes_slide creates it (has_notes_slide becomes True) even
        # though nothing is ever written into it.
        slide.notes_slide.notes_text_frame.text = "   "
        prs.save(source)

        context = make_context(source, tmp_path)
        convert_with_markitdown(context)
        assert "### Notes:" not in context.markdown_output

    def test_pptx_real_notes_are_kept(self, tmp_path: Path):
        from pptx import Presentation

        source = tmp_path / "notes.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide with real notes"
        slide.notes_slide.notes_text_frame.text = "Remember to mention the deadline."
        prs.save(source)

        context = make_context(source, tmp_path)
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "### Notes:" in md
        assert "Remember to mention the deadline." in md

    def test_pptx_literal_bullet_glyphs_become_list_items(self, tmp_path: Path):
        """A slide body sometimes has a literal "•" typed into the text
        instead of PowerPoint's own bullet-list paragraph formatting (common
        in decks pasted in from elsewhere) -- must normalize the same way a
        PDF's bullets do, not leak the raw glyph."""
        from pptx import Presentation
        from pptx.util import Inches

        source = tmp_path / "bullets.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
        tf = box.text_frame
        tf.text = "• first point"
        tf.add_paragraph().text = "• second point"
        prs.save(source)

        context = make_context(source, tmp_path)
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "- first point" in md
        assert "- second point" in md
        assert not any(line and line[0] in "●•" for line in md.splitlines())

    def test_pptx_notes_bullets_are_also_normalized(self, tmp_path: Path):
        """Same normalization applies to presenter notes text, not only the
        visible slide body -- confirmed against a real deck where this
        content lived in the notes, not on any slide."""
        from pptx import Presentation

        source = tmp_path / "notes-bullets.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Methodology"
        slide.notes_slide.notes_text_frame.text = (
            "Summary of choices:\n• point one\n• point two"
        )
        prs.save(source)

        context = make_context(source, tmp_path)
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "- point one" in md
        assert "- point two" in md
        assert not any(line and line[0] in "●•" for line in md.splitlines())

    def test_pptx_placeholder_alt_text_is_not_leaked(self, tmp_path: Path):
        """A filename Google Slides wrote into a shape's description (e.g.
        because the author never set real alt text) must not be shown
        verbatim -- it reads as a bug, not a description, to a reader."""
        from io import BytesIO

        from PIL import Image
        from pptx import Presentation

        source = tmp_path / "alt.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        buffer = BytesIO()
        Image.new("RGB", (400, 300), (40, 120, 200)).save(buffer, format="PNG")
        buffer.seek(0)
        picture = slide.shapes.add_picture(buffer, 0, 0)
        picture._element._nvXxPr.cNvPr.attrib["descr"] = "preencoded.png"
        prs.save(source)

        context = make_context(source, tmp_path)
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "preencoded" not in md
        assert "Slide image" in md

    def test_pptx_real_alt_text_is_kept(self, tmp_path: Path):
        from io import BytesIO

        from PIL import Image
        from pptx import Presentation

        source = tmp_path / "alt.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        buffer = BytesIO()
        Image.new("RGB", (400, 300), (40, 120, 200)).save(buffer, format="PNG")
        buffer.seek(0)
        picture = slide.shapes.add_picture(buffer, 0, 0)
        picture._element._nvXxPr.cNvPr.attrib["descr"] = "Diagram of the water cycle"
        prs.save(source)

        context = make_context(source, tmp_path)
        convert_with_markitdown(context)
        assert "Diagram of the water cycle" in context.markdown_output

    def test_stats_count_pages_slides_sheets(self, tmp_path: Path):
        pdf = make_deck_pdf(tmp_path / "deck.pdf")
        pdf_context = make_context(pdf, tmp_path, ocr_mode="auto")
        pdf_doc = convert_with_markitdown(pdf_context)
        assert pdf_doc.stats.pages == 4

        pptx = make_pptx(tmp_path / "slides.pptx", slides=3)
        pptx_context = make_context(pptx, tmp_path)
        pptx_doc = convert_with_markitdown(pptx_context)
        assert pptx_doc.stats.slides == 3

        xlsx = make_xlsx(tmp_path / "book.xlsx", sheets=2)
        xlsx_context = make_context(xlsx, tmp_path)
        xlsx_doc = convert_with_markitdown(xlsx_context)
        assert xlsx_doc.stats.sheets == 2

    def test_ocr_never_still_marks_textless_pages(self, tmp_path: Path):
        source = make_scanned_pdf(tmp_path / "scan.pdf", pages=1)
        context = make_context(source, tmp_path, ocr_mode="never")
        doc = convert_with_markitdown(context)
        assert not context.ocr_used
        assert "[Image page — no text]" in context.markdown_output
        assert not any(w.code == "ocr_unavailable" for w in doc.warnings)

    def test_ocr_unavailable_warning(self, tmp_path: Path, monkeypatch):
        import converters.markitdown as mod

        monkeypatch.setattr(mod, "tesseract_available", lambda: False)
        source = make_scanned_pdf(tmp_path / "scan.pdf")
        context = make_context(source, tmp_path, ocr_mode="auto")
        doc = convert_with_markitdown(context)
        codes = [w.code for w in doc.warnings]
        assert "ocr_unavailable" in codes
        assert not context.ocr_used

    def test_ocr_unavailable_warning_when_package_missing(self, tmp_path: Path, monkeypatch):
        import converters.markitdown as mod

        monkeypatch.setattr(mod, "tesseract_available", lambda: True)
        monkeypatch.setattr(mod, "_ocrmypdf_importable", lambda: False)
        source = make_scanned_pdf(tmp_path / "scan-missing-pkg.pdf")
        context = make_context(source, tmp_path, ocr_mode="auto")
        doc = convert_with_markitdown(context)
        codes = [w.code for w in doc.warnings]
        assert "ocr_unavailable" in codes
        assert not context.ocr_used

    def test_ocr_prepass_runs_when_available(self, tmp_path: Path):
        from app.services.ocr import tesseract_available

        if not tesseract_available():
            pytest.skip("tesseract not installed on this system")

        source = make_scanned_pdf(tmp_path / "scan2.pdf")
        context = make_context(source, tmp_path, ocr_mode="auto")
        doc = convert_with_markitdown(context)
        codes = [w.code for w in doc.warnings]
        assert "ocr_used" in codes
        assert context.ocr_used
        assert doc.stats.ocr_pages >= 1

    def test_ocr_never_skips_prepass(self, tmp_path: Path, monkeypatch):
        import converters.markitdown as mod

        monkeypatch.setattr(mod, "tesseract_available", lambda: True)
        source = make_scanned_pdf(tmp_path / "scan3.pdf")
        context = make_context(source, tmp_path, ocr_mode="never")
        doc = convert_with_markitdown(context)
        assert not context.ocr_used
        assert not any(w.code == "ocr_unavailable" for w in doc.warnings)

    def test_ocr_always_forces_prepass_without_detection(self, tmp_path: Path, monkeypatch):
        import converters.markitdown as mod

        calls: list[bool] = []

        def fake_detect(source) -> bool:
            raise AssertionError("detection must be bypassed in 'always' mode")

        def fake_ocr(source, context, *, force_ocr):
            calls.append(force_ocr)
            return None

        monkeypatch.setattr(mod, "_detect_scanned_pdf", fake_detect)
        monkeypatch.setattr(mod, "tesseract_available", lambda: True)
        monkeypatch.setattr(mod, "_ocr_pdf_copy", fake_ocr)

        source = make_pdf(tmp_path / "plain.pdf", pages=2)
        context = make_context(source, tmp_path, ocr_mode="always")
        convert_with_markitdown(context)
        assert calls == [True]

    def test_ocr_auto_skips_prepass_when_not_scanned(self, tmp_path: Path, monkeypatch):
        import converters.markitdown as mod

        calls: list[bool] = []

        def fake_ocr(source, context, *, force_ocr):
            calls.append(force_ocr)
            return None

        monkeypatch.setattr(mod, "_detect_scanned_pdf", lambda source: False)
        monkeypatch.setattr(mod, "tesseract_available", lambda: True)
        monkeypatch.setattr(mod, "_ocr_pdf_copy", fake_ocr)

        source = make_pdf(tmp_path / "plain2.pdf", pages=2)
        context = make_context(source, tmp_path, ocr_mode="auto")
        convert_with_markitdown(context)
        assert calls == []
        assert not context.ocr_used


class TestColumnAwarePdf:
    def test_columns_not_interleaved(self, tmp_path: Path):
        source = make_two_column_pdf(tmp_path / "cols.pdf")
        context = make_context(source, tmp_path, ocr_mode="never")
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "Left column line A\nLeft column line B\nLeft column line C" in md
        assert md.index("Left column line A") < md.index("Right column line A")

    def test_extraction_decomposes_ligatures(self, tmp_path: Path, monkeypatch):
        """Typographic ligatures (ﬀ, ﬁ, ﬂ, ﬃ) must come out as plain letters.

        Fonts embedded by real PDF export toolchains (PowerPoint, Google
        Slides, LibreOffice) commonly substitute "ff"/"fi"/"fl" pairs with a
        single ligature glyph. PyMuPDF preserves that glyph by default, so
        without the TEXT_PRESERVE_LIGATURES flag cleared, "Effects" extracts
        as "Eﬀects" and silently fails a reader's search or copy-paste.
        A base14 test font has no ligature glyphs to reproduce this with, so
        this asserts the flag itself reaches PyMuPDF rather than depending on
        font substitution.
        """
        import pymupdf

        from converters.markitdown_pdf import _TEXT_DICT_FLAGS

        assert not (_TEXT_DICT_FLAGS & pymupdf.TEXT_PRESERVE_LIGATURES)

        source = make_pdf(tmp_path / "flags.pdf", pages=1)
        seen_flags: list[int] = []
        original = pymupdf.Page.get_text

        def spy(self, *args, **kwargs):
            if args and args[0] == "dict":
                # PyMuPDF defaults to TEXTFLAGS_DICT (ligatures preserved) when
                # no ``flags`` kwarg is passed -- that missing kwarg is exactly
                # what regressed here, so it must count as "ligatures on",
                # not as an empty/zero flag set.
                seen_flags.append(kwargs.get("flags", pymupdf.TEXTFLAGS_DICT))
            return original(self, *args, **kwargs)

        monkeypatch.setattr(pymupdf.Page, "get_text", spy)
        context = make_context(source, tmp_path, ocr_mode="never")
        convert_with_markitdown(context)

        assert seen_flags, "the page.get_text('dict', ...) call was not exercised"
        assert all(not (f & pymupdf.TEXT_PRESERVE_LIGATURES) for f in seen_flags)

    def test_wrapped_sentence_reflows_into_one_line(self, tmp_path: Path):
        """A sentence PDF-wrapped across several same-size lines becomes one."""
        import pymupdf

        source = tmp_path / "wrapped.pdf"
        doc = pymupdf.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Section heading", fontsize=20)
        y = 120
        for chunk in (
            "the quick brown fox",
            "jumps over the lazy",
            "dog while the sun",
            "was setting slowly.",
        ):
            page.insert_text((72, y), chunk, fontsize=12)
            y += 20
        doc.save(source)
        doc.close()

        context = make_context(source, tmp_path, ocr_mode="never")
        convert_with_markitdown(context)
        md = context.markdown_output
        assert (
            "the quick brown fox jumps over the lazy dog while the sun "
            "was setting slowly."
        ) in md
        # It must have become exactly one line, not four.
        assert "the quick brown fox\n" not in md

    def test_capitalized_label_does_not_merge_into_paragraph(self, tmp_path: Path):
        """A new capitalized line must not glue onto unpunctuated prose above it.

        PDFs that flatten a table into loose lines (a row label like "Meaning"
        followed by its cell text) must not have the label silently absorbed
        into the previous line just because that line lacks a period.
        """
        import pymupdf

        source = tmp_path / "label.pdf"
        doc = pymupdf.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Overview", fontsize=20)
        page.insert_text(
            (72, 120), "some introductory text without a period", fontsize=12
        )
        page.insert_text((72, 140), "Meaning", fontsize=12)
        doc.save(source)
        doc.close()

        context = make_context(source, tmp_path, ocr_mode="never")
        convert_with_markitdown(context)
        lines = context.markdown_output.splitlines()
        assert "some introductory text without a period" in lines
        assert "Meaning" in lines

    def test_bullet_does_not_merge_into_preceding_paragraph(self, tmp_path: Path):
        """A standalone marker line ("-") followed by content must start a new
        list item, never get glued onto unpunctuated prose above it.

        Uses a plain ASCII "-" rather than a symbol glyph like "●": PyMuPDF's
        built-in base14 fonts don't carry symbol glyphs and silently
        substitute a different character on insertion (a synthetic-PDF-fixture
        limitation, not a property of real embedded fonts -- verified directly
        against a real PDF that "●" round-trips correctly in production).
        """
        import pymupdf

        source = tmp_path / "bullet.pdf"
        doc = pymupdf.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Overview", fontsize=20)
        page.insert_text(
            (72, 120), "introductory text with no terminal punctuation", fontsize=12
        )
        page.insert_text((72, 140), "-", fontsize=12)
        page.insert_text((72, 160), "first bullet point", fontsize=12)
        page.insert_text((72, 180), "-", fontsize=12)
        page.insert_text((72, 200), "second bullet point", fontsize=12)
        doc.save(source)
        doc.close()

        context = make_context(source, tmp_path, ocr_mode="never")
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "introductory text with no terminal punctuation" in md
        assert "- first bullet point" in md
        assert "- second bullet point" in md
        assert "punctuation - first" not in md

    def test_bullet_glyphs_become_markdown_list_items(self, tmp_path: Path):
        """Wiring check: `_extract_page` actually calls the normalizer, for
        both the marker-merge branch and the inline plain-content branch.
        Uses portable ASCII markers only -- see the note on the previous test
        for why symbol glyphs aren't reliable in a synthetic test PDF; full
        glyph-set coverage is unit-tested directly against
        `_normalize_bullet_prefix` in TestNormalizeBulletPrefix below.
        """
        import pymupdf

        source = tmp_path / "list.pdf"
        doc = pymupdf.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Checklist", fontsize=20)
        # Inline marker already on the same line as its text, as it arrives
        # for the unambiguous glyphs (●, ✔) in a real PDF -- a dash needs the
        # trailing space (see TestNormalizeBulletPrefix) to stay distinct
        # from a hyphenated word, so it's included here with one.
        page.insert_text((72, 120), "- inline item", fontsize=12)
        # A standalone marker line followed by a separate text line -- the
        # marker-merge path used when a PDF renders the icon as its own run.
        page.insert_text((72, 140), "-", fontsize=12)
        page.insert_text((72, 160), "standalone marker item", fontsize=12)
        page.insert_text((72, 180), "1. numbered item", fontsize=12)
        doc.save(source)
        doc.close()

        context = make_context(source, tmp_path, ocr_mode="never")
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "- inline item" in md
        assert "- standalone marker item" in md
        assert "1. numbered item" in md

    def test_repeated_standalone_markers_collapse_to_one(self, tmp_path: Path):
        """Several identical standalone marker lines in a row (observed in a
        real document: four consecutive "-" lines used as a divider before a
        section label) are one decorative divider, not four list items --
        must not leave the extras stuck mid-line as literal characters.
        """
        import pymupdf

        source = tmp_path / "divider.pdf"
        doc = pymupdf.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Reference", fontsize=20)
        for y in (120, 140, 160, 180):
            page.insert_text((72, y), "-", fontsize=12)
        page.insert_text((72, 200), "Section label", fontsize=12)
        doc.save(source)
        doc.close()

        context = make_context(source, tmp_path, ocr_mode="never")
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "- Section label" in md
        assert "- - " not in md
        assert md.count("Section label") == 1


class TestNormalizeBulletPrefix:
    """Direct unit tests for the glyph-normalization helper.

    Kept independent of PDF rendering: PyMuPDF's built-in base14 fonts don't
    carry every symbol glyph used here (confirmed empirically -- inserting
    "●" via the default font round-trips as a different character), which is
    a limitation of synthesizing a *test* PDF, not of extracting from a real
    one. Testing the pure function directly is both more reliable and a more
    honest match for what's actually being verified.
    """

    def test_bullet_glyphs(self):
        from converters.markitdown_pdf import _normalize_bullet_prefix

        for glyph in "●•◦▪‣":
            assert _normalize_bullet_prefix(f"{glyph}item text") == "- item text"
            assert _normalize_bullet_prefix(f"{glyph} item text") == "- item text"

    def test_check_glyphs(self):
        from converters.markitdown_pdf import _normalize_bullet_prefix

        for glyph in "✔✓☑☒":
            assert (
                _normalize_bullet_prefix(f"{glyph}done") == "- ✔ done"
            )

    def test_numbered_marker(self):
        from converters.markitdown_pdf import _normalize_bullet_prefix

        assert _normalize_bullet_prefix("1. first item") == "1. first item"
        assert _normalize_bullet_prefix("2) second item") == "2. second item"

    def test_dash_requires_trailing_space(self):
        from converters.markitdown_pdf import _normalize_bullet_prefix

        assert _normalize_bullet_prefix("- dash item") == "- dash item"
        # No space after the dash -- likely a hyphenated word, not a bullet.
        assert _normalize_bullet_prefix("-verified") == "-verified"

    def test_decimal_number_is_untouched(self):
        from converters.markitdown_pdf import _normalize_bullet_prefix

        assert (
            _normalize_bullet_prefix("3.5 million students")
            == "3.5 million students"
        )

    def test_negative_number_is_untouched(self):
        from converters.markitdown_pdf import _normalize_bullet_prefix

        assert _normalize_bullet_prefix("-5 degrees") == "-5 degrees"

    def test_plain_text_is_untouched(self):
        from converters.markitdown_pdf import _normalize_bullet_prefix

        assert _normalize_bullet_prefix("ordinary sentence.") == "ordinary sentence."

    def test_individually_bold_wrapped_glyph_is_still_recognized(self):
        """`_wrap_bold` can wrap just the bullet span (not the whole line)
        when only the glyph was bold in the source PDF, producing "**●**
        text" rather than "●text". Confirmed in real converted output."""
        from converters.markitdown_pdf import _normalize_bullet_prefix

        assert (
            _normalize_bullet_prefix("**●** Log in to your account")
            == "- Log in to your account"
        )
        assert _normalize_bullet_prefix("**✔** done") == "- ✔ done"
        assert _normalize_bullet_prefix("**1.** first item") == "1. first item"

    def test_decimal_number_is_not_mistaken_for_a_list_marker(self, tmp_path: Path):
        import pymupdf

        source = tmp_path / "decimal.pdf"
        doc = pymupdf.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Report", fontsize=20)
        page.insert_text((72, 120), "3.5 million students were affected.", fontsize=12)
        doc.save(source)
        doc.close()

        context = make_context(source, tmp_path, ocr_mode="never")
        convert_with_markitdown(context)
        assert "3.5 million students were affected." in context.markdown_output

    def test_large_font_promoted_to_heading(self, tmp_path: Path):
        source = make_pdf(tmp_path / "headings.pdf", pages=2)
        context = make_context(source, tmp_path, ocr_mode="never")
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "## Chapter 1" in md
        assert "## Paragraph text on page 1." not in md

    def test_page_markers_kept_in_fidelity(self, tmp_path: Path):
        source = make_pdf(tmp_path / "markers.pdf", pages=2)
        context = make_context(source, tmp_path, ocr_mode="never", output_mode="fidelity")
        convert_with_markitdown(context)
        assert "<!-- Page 1 -->" in context.markdown_output
        assert "<!-- Page 2 -->" in context.markdown_output

    def test_page_markers_stripped_in_clean(self, tmp_path: Path):
        source = make_pdf(tmp_path / "markers2.pdf", pages=2)
        context = make_context(source, tmp_path, ocr_mode="never", output_mode="clean")
        convert_with_markitdown(context)
        assert "<!-- Page" not in context.markdown_output

    def test_checkmark_merged_with_following_line(self, tmp_path: Path):
        import os

        import pymupdf

        symbol_font = r"C:\Windows\Fonts\seguisym.ttf"
        if not os.path.exists(symbol_font):
            pytest.skip("Segoe UI Symbol font not available")

        source = tmp_path / "checks.pdf"
        doc = pymupdf.open()
        page = doc.new_page()
        page.insert_font(fontname="f-sym", fontfile=symbol_font)
        page.insert_text((72, 72), "Task list", fontsize=14)
        page.insert_text((72, 110), "✔", fontname="f-sym", fontsize=12)
        page.insert_text((72, 130), "Understand the material", fontsize=12)
        doc.save(source)
        doc.close()
        context = make_context(source, tmp_path, ocr_mode="never")
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "✔ Understand the material" in md

    def test_standalone_bullet_and_number_merged(self, tmp_path: Path):
        import pymupdf

        source = tmp_path / "list.pdf"
        doc = pymupdf.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Steps", fontsize=14)
        page.insert_text((72, 110), "1.", fontsize=12)
        page.insert_text((72, 132), "Humanap ng taong pinagkakatiwalaan.", fontsize=12)
        page.insert_text((72, 160), "2.", fontsize=12)
        page.insert_text((72, 182), "I-dokumento ang insidente.", fontsize=12)
        doc.save(source)
        doc.close()
        context = make_context(source, tmp_path, ocr_mode="never")
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "1. Humanap ng taong pinagkakatiwalaan." in md
        assert "2. I-dokumento ang insidente." in md
        assert "\n1.\n" not in md

    def test_placeholder_single_brackets(self, tmp_path: Path):
        source = make_deck_pdf(tmp_path / "deck.pdf")
        context = make_context(source, tmp_path, ocr_mode="never")
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "[Slide image — no text]" in md
        assert "[[Slide image — no text]]" not in md


class TestDuplicatePages:
    def test_duplicates_removed_with_warning(self, tmp_path: Path):
        source = make_deck_pdf(tmp_path / "deck.pdf")
        context = make_context(source, tmp_path, ocr_mode="never")
        doc = convert_with_markitdown(context)
        codes = [w.code for w in doc.warnings]
        assert "duplicate_pages_removed" in codes
        md = context.markdown_output
        assert md.count("Mission 3") == 1

    def test_identical_textless_pages_not_deduplicated(self, tmp_path: Path):
        from tests.fixtures.make_fixtures import make_scanned_pdf

        source = make_scanned_pdf(tmp_path / "scans.pdf", pages=3)
        context = make_context(source, tmp_path, ocr_mode="never")
        doc = convert_with_markitdown(context)
        assert not any(w.code == "duplicate_pages_removed" for w in doc.warnings)


class TestDeckVsScanOcr:
    def test_deck_mode_skips_ocr_with_placeholder(self, tmp_path: Path, monkeypatch):
        import converters.markitdown as mod

        def fail_ocr(*args, **kwargs):
            raise AssertionError("OCR must not run in deck mode")

        monkeypatch.setattr(mod, "_ocr_pdf_copy", fail_ocr)
        source = make_deck_pdf(tmp_path / "deck2.pdf")
        context = make_context(source, tmp_path, ocr_mode="auto")
        doc = convert_with_markitdown(context)
        assert not context.ocr_used
        codes = [w.code for w in doc.warnings]
        assert "decorative_pages_skipped" in codes
        assert "[Slide image — no text]" in context.markdown_output

    def test_scan_mode_ocrs_all_pages(self, tmp_path: Path, monkeypatch):
        import converters.markitdown as mod

        calls: list[bool] = []

        def record_ocr(source, context, *, force_ocr):
            calls.append(force_ocr)
            return None

        monkeypatch.setattr(mod, "_ocr_pdf_copy", record_ocr)
        monkeypatch.setattr(mod, "tesseract_available", lambda: True)
        source = make_scanned_pdf(tmp_path / "scan4.pdf", pages=2)
        context = make_context(source, tmp_path, ocr_mode="auto")
        convert_with_markitdown(context)
        assert calls == [False]
        assert not context.ocr_used


class TestPptxHeadings:
    def test_large_font_text_box_promoted(self, tmp_path: Path):
        from pptx import Presentation
        from pptx.util import Pt

        source = tmp_path / "deck.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(0, 0, 9_000_000, 1_000_000)
        title_box.text_frame.text = "Big Mission Title"
        title_box.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
        body_box = slide.shapes.add_textbox(0, 2_000_000, 9_000_000, 3_000_000)
        body_box.text_frame.text = "Normal body text"
        body_box.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
        prs.save(source)

        context = make_context(source, tmp_path)
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "## Big Mission Title" in md
        assert "## Normal body text" not in md

    def test_title_placeholder_still_promoted(self, tmp_path: Path):
        source = make_pptx(tmp_path / "deck2.pptx", slides=2)
        context = make_context(source, tmp_path)
        convert_with_markitdown(context)
        md = context.markdown_output
        assert "# Introduction" in md

def _write(path: Path, data: bytes | str) -> Path:
    path.write_bytes(data.encode("utf-8") if isinstance(data, str) else data)
    return path


def make_epub(path: Path) -> Path:
    """Minimal but valid EPUB 3 package."""
    import zipfile

    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("mimetype", "application/epub+zip")
        archive.writestr(
            "META-INF/container.xml",
            '<?xml version="1.0"?><container version="1.0" '
            'xmlns="urn:oasis:names:tc:opendocument:xmlns:container"><rootfiles>'
            '<rootfile full-path="OEBPS/content.opf" '
            'media-type="application/oebps-package+xml"/></rootfiles></container>',
        )
        archive.writestr(
            "OEBPS/content.opf",
            '<?xml version="1.0"?><package xmlns="http://www.idpf.org/2007/opf" '
            'version="3.0" unique-identifier="id"><metadata '
            'xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:title>Test Book</dc:title>'
            "<dc:creator>An Author</dc:creator>"
            '<dc:identifier id="id">urn:uuid:1</dc:identifier>'
            "<dc:language>en</dc:language></metadata><manifest>"
            '<item id="c1" href="c1.xhtml" media-type="application/xhtml+xml"/>'
            '</manifest><spine><itemref idref="c1"/></spine></package>',
        )
        archive.writestr(
            "OEBPS/c1.xhtml",
            '<?xml version="1.0"?><html xmlns="http://www.w3.org/1999/xhtml"><body>'
            "<h1>Chapter One</h1><p>Body text of the chapter.</p></body></html>",
        )
    return path


class TestExpandedFormats:
    """Every extension in ALLOWED_EXTENSIONS must actually convert.

    The allowlist and the registered converters are two halves of one contract;
    widening one without the other fails at conversion time instead of upload.
    """

    @pytest.mark.parametrize(
        ("name", "content", "expected"),
        [
            ("page.html", "<html><body><h1>Title</h1><p>Body</p></body></html>", "# Title"),
            ("page.htm", "<html><body><h2>Sub</h2></body></html>", "## Sub"),
            ("notes.txt", "plain text line\nsecond line\n", "plain text line"),
            ("notes.md", "# Heading\n\n- a\n- b\n", "# Heading"),
            ("data.csv", "name,score\nAda,99\nGrace,100\n", "| Ada | 99 |"),
            ("data.tsv", "name\tscore\nAda\t99\n", "Ada"),
            ("data.json", '{"a": 1, "b": [1, 2]}', '"a"'),
            ("data.xml", "<?xml version='1.0'?><root><item>x</item></root>", "item"),
        ],
    )
    def test_text_family_converts(self, tmp_path: Path, name, content, expected):
        source = _write(tmp_path / name, content)
        context = make_context(source, tmp_path)
        doc = convert_with_markitdown(context)
        assert isinstance(doc, Document)
        assert expected in context.markdown_output

    def test_ipynb_converts_cells(self, tmp_path: Path):
        import json

        notebook = {
            "cells": [
                {"cell_type": "markdown", "metadata": {}, "source": ["# NB Title"]},
                {
                    "cell_type": "code",
                    "execution_count": 1,
                    "metadata": {},
                    "outputs": [],
                    "source": ["print('hi')"],
                },
            ],
            "metadata": {
                "kernelspec": {
                    "display_name": "Python 3",
                    "language": "python",
                    "name": "python3",
                }
            },
            "nbformat": 4,
            "nbformat_minor": 5,
        }
        source = _write(tmp_path / "nb.ipynb", json.dumps(notebook))
        context = make_context(source, tmp_path)
        convert_with_markitdown(context)
        assert "# NB Title" in context.markdown_output
        assert "print('hi')" in context.markdown_output

    def test_epub_converts_with_metadata(self, tmp_path: Path):
        source = make_epub(tmp_path / "book.epub")
        context = make_context(source, tmp_path)
        convert_with_markitdown(context)
        assert "Test Book" in context.markdown_output
        assert "Chapter One" in context.markdown_output


def _png(width: int, height: int, color=(40, 120, 200)) -> bytes:
    from io import BytesIO

    from PIL import Image, ImageDraw

    image = Image.new("RGB", (width, height), color)
    # A little internal structure so the file does not compress to nothing and
    # trip the MIN_BYTES branch of the decorative filter.
    ImageDraw.Draw(image).ellipse(
        (width * 0.1, height * 0.1, width * 0.9, height * 0.9), fill=(250, 210, 70)
    )
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def make_figure_pdf(path: Path, pages: int = 2, decorations: bool = True) -> Path:
    """A PDF with one real figure per page, plus decorative noise."""
    import pymupdf

    doc = pymupdf.open()
    for number in range(1, pages + 1):
        page = doc.new_page()
        page.insert_text((72, 72), f"Chapter {number}: findings and discussion")
        page.insert_image(
            pymupdf.Rect(72, 110, 400, 350),
            stream=_png(500, 360, (20, 80 + number * 30, 190)),
        )
        if decorations:
            # A hairline rule and a spacer pixel: both must be filtered out.
            page.insert_image(
                pymupdf.Rect(72, 700, 500, 703), stream=_png(600, 4, (10, 10, 10))
            )
            page.insert_image(
                pymupdf.Rect(520, 40, 523, 43), stream=_png(3, 3, (10, 10, 10))
            )
    doc.save(path)
    doc.close()
    return path


class TestImageExtraction:
    def test_pdf_figures_are_anchored_to_their_page(self, tmp_path: Path):
        source = make_figure_pdf(tmp_path / "figures.pdf", pages=3)
        context = make_context(source, tmp_path)
        doc = convert_with_markitdown(context)

        assert doc.stats.images == 3
        out = context.markdown_output
        for page_number in (1, 2, 3):
            marker = f"<!-- Page {page_number} -->"
            after = out.index(marker) + len(marker)
            # The figure must follow its own page marker, before the next one.
            nxt = out.find("<!-- Page ", after)
            segment = out[after:] if nxt == -1 else out[after:nxt]
            assert "](assets/" in segment, f"page {page_number} lost its figure"

    def test_decorative_images_are_filtered_out(self, tmp_path: Path):
        source = make_figure_pdf(tmp_path / "noisy.pdf", pages=2)
        context = make_context(source, tmp_path)
        doc = convert_with_markitdown(context)
        # Two real figures; the rules and spacer pixels must not survive.
        assert doc.stats.images == 2

    def test_assets_are_written_to_the_output_directory(self, tmp_path: Path):
        source = make_figure_pdf(tmp_path / "figures.pdf", pages=2)
        context = make_context(source, tmp_path)
        convert_with_markitdown(context)

        assets = sorted((context.output_dir / "assets").glob("*"))
        assert [a.name for a in assets] == ["image-001.png", "image-002.png"]
        assert all(a.stat().st_size > 0 for a in assets)

    def test_repeated_image_is_saved_once(self, tmp_path: Path):
        """A logo on every page should cost one file, not one per page."""
        import pymupdf

        logo = _png(300, 200, (90, 90, 200))
        source = tmp_path / "logo.pdf"
        doc = pymupdf.open()
        for number in range(1, 5):
            page = doc.new_page()
            # Unique text per page, or the duplicate-page pre-pass collapses
            # them and this stops testing image dedup at all.
            page.insert_text((72, 72), f"Section {number} body text")
            page.insert_image(pymupdf.Rect(72, 100, 272, 233), stream=logo)
        doc.save(source)
        doc.close()

        context = make_context(source, tmp_path)
        convert_with_markitdown(context)
        assets = list((context.output_dir / "assets").glob("*"))
        assert len(assets) == 1

    def test_background_on_every_page_is_referenced_once(self, tmp_path: Path):
        """A logo on every page is page-template chrome, not four figures.

        save_image() already writes it to disk once (test_repeated_image_is_
        saved_once). This is the separate step: an image on all 4 of 4 pages
        (100% > the 50% threshold) should be *referenced* once too, kept on
        its first page, not repeated in the Markdown on every page.
        """
        import pymupdf

        logo = _png(300, 200, (90, 90, 200))
        source = tmp_path / "logo.pdf"
        doc = pymupdf.open()
        for number in range(1, 5):
            page = doc.new_page()
            page.insert_text((72, 72), f"Section {number} body text")
            page.insert_image(pymupdf.Rect(72, 100, 272, 233), stream=logo)
        doc.save(source)
        doc.close()

        context = make_context(source, tmp_path)
        result = convert_with_markitdown(context)
        assert context.markdown_output.count("](assets/") == 1
        assert "<!-- Page 1 -->" in context.markdown_output
        page_1_segment = context.markdown_output.split("<!-- Page 2 -->")[0]
        assert "](assets/" in page_1_segment
        assert any(
            w.code == "recurring_backgrounds_suppressed" for w in result.warnings
        )

    def test_image_on_a_minority_of_pages_is_not_suppressed(self, tmp_path: Path):
        """A figure that only recurs a couple of times must survive intact."""
        import pymupdf

        chart = _png(400, 300, (40, 120, 200))
        source = tmp_path / "chart.pdf"
        doc = pymupdf.open()
        for number in range(1, 7):
            page = doc.new_page()
            page.insert_text((72, 72), f"Section {number} body text")
            # Appears on only 2 of 6 pages (33%), well under the threshold.
            if number in (1, 4):
                page.insert_image(pymupdf.Rect(72, 100, 372, 250), stream=chart)
        doc.save(source)
        doc.close()

        context = make_context(source, tmp_path)
        result = convert_with_markitdown(context)
        assert context.markdown_output.count("](assets/") == 2
        assert not any(
            w.code == "recurring_backgrounds_suppressed" for w in result.warnings
        )

    def test_extraction_can_be_turned_off(self, tmp_path: Path):
        source = make_figure_pdf(tmp_path / "figures.pdf", pages=2)
        context = make_context(source, tmp_path, extract_images=False)
        doc = convert_with_markitdown(context)

        assert doc.stats.images == 0
        assert "](assets/" not in context.markdown_output
        assert not (context.output_dir / "assets").exists()

    def test_figures_survive_preserve_links_off(self, tmp_path: Path):
        """Turning off links must not silently delete the document's figures."""
        source = make_figure_pdf(tmp_path / "figures.pdf", pages=2)
        context = make_context(source, tmp_path, preserve_links=False)
        doc = convert_with_markitdown(context)
        assert doc.stats.images == 2

    def test_extraction_reports_a_warning(self, tmp_path: Path):
        source = make_figure_pdf(tmp_path / "figures.pdf", pages=2)
        context = make_context(source, tmp_path)
        doc = convert_with_markitdown(context)
        assert any(w.code == "images_extracted" for w in doc.warnings)

    def test_docx_figures_are_placed_inline(self, tmp_path: Path):
        from io import BytesIO

        import docx
        from docx.shared import Inches

        source = tmp_path / "report.docx"
        document = docx.Document()
        document.add_heading("Report", level=1)
        document.add_paragraph("Before the figure.")
        document.add_picture(BytesIO(_png(520, 340)), width=Inches(4))
        document.add_paragraph("Between the figures.")
        document.add_picture(BytesIO(_png(500, 320, (40, 160, 90))), width=Inches(4))
        document.add_paragraph("After the figure.")
        document.save(source)

        context = make_context(source, tmp_path)
        doc = convert_with_markitdown(context)
        out = context.markdown_output

        assert doc.stats.images == 2
        # Inline placement: each figure sits between its neighbouring paragraphs.
        assert out.index("Before the figure.") < out.index("](assets/")
        assert out.index("](assets/") < out.index("Between the figures.")
        assert out.index("Between the figures.") < out.rindex("](assets/")
        assert out.rindex("](assets/") < out.index("After the figure.")
        # The stub MarkItDown leaves behind must be fully consumed.
        assert "data:image" not in out

    def test_pptx_picture_is_saved_and_referenced(self, tmp_path: Path):
        from io import BytesIO

        from pptx import Presentation
        from pptx.util import Inches

        source = tmp_path / "deck.pptx"
        prs = Presentation()
        first = prs.slides.add_slide(prs.slide_layouts[5])
        first.shapes.title.text = "Deck"
        second = prs.slides.add_slide(prs.slide_layouts[5])
        second.shapes.title.text = "Architecture"
        second.shapes.add_picture(
            BytesIO(_png(480, 320)), Inches(1), Inches(2), Inches(4), Inches(2.6)
        )
        # A spacer that must leave nothing at all behind.
        second.shapes.add_picture(
            BytesIO(_png(6, 6, (0, 0, 0))),
            Inches(6),
            Inches(0.4),
            Inches(0.1),
            Inches(0.1),
        )
        prs.save(source)

        context = make_context(source, tmp_path)
        doc = convert_with_markitdown(context)
        out = context.markdown_output

        assert doc.stats.images == 1
        assert "](assets/" in out
        assert "[Image:" not in out
        # The figure belongs to slide 2, not slide 1.
        assert out.index("<!-- Slide number: 2 -->") < out.index("](assets/")


class TestMergeAwareXlsx:
    """The stock pandas-based xlsx converter has three confirmed defects:
    an empty cell renders as the literal string "NaN", an empty sheet
    produces a malformed one-cell table, and a merged cell's non-top-left
    columns silently vanish from the table's column count. These test the
    replacement converter directly (bypassing pandas) against real
    openpyxl-written and -read files.
    """

    def _convert(self, tmp_path: Path, workbook) -> str:
        source = tmp_path / "book.xlsx"
        workbook.save(source)
        context = make_context(source, tmp_path)
        convert_with_markitdown(context)
        return context.markdown_output

    def test_empty_cell_is_blank_not_nan(self, tmp_path: Path):
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Grades"
        ws.append(["Student", "Score", "Notes"])
        ws.append(["Ada", 99, None])
        ws.append(["Grace", 100, "top"])

        md = self._convert(tmp_path, wb)
        assert "NaN" not in md
        assert "| Ada | 99 |  |" in md

    def test_empty_sheet_gets_a_clear_note_not_a_malformed_table(
        self, tmp_path: Path
    ):
        import openpyxl

        wb = openpyxl.Workbook()
        wb.active.title = "Empty"

        md = self._convert(tmp_path, wb)
        assert "_(empty sheet)_" in md
        assert "|" not in md

    def test_merged_cell_repeats_its_value_instead_of_collapsing_the_row(
        self, tmp_path: Path
    ):
        """A merged A1:B1 cell's non-top-left value is genuinely gone the
        moment the file is saved -- confirmed directly against the raw XLSX
        XML, which stores only one <c> element for the whole merged range.
        There is nothing left to "recover"; the defect being fixed is that
        the stock converter rendered that row with *fewer columns* than the
        table's header, a malformed/ragged table. The fix keeps every row's
        column count consistent by repeating the one value that survives.
        """
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Merged"
        ws.append(["Q1", "Q2"])
        ws.merge_cells("A1:B1")

        md = self._convert(tmp_path, wb)
        assert "| Q1 | Q1 |" in md
        assert "Q2" not in md

    def test_whole_number_is_not_shown_as_a_float(self, tmp_path: Path):
        from converters.markitdown_xlsx import _format_cell_value

        assert _format_cell_value(99.0) == "99"
        assert _format_cell_value(99.5) == "99.5"
        assert _format_cell_value(None) == ""
        assert _format_cell_value("") == ""
        assert _format_cell_value("Ada") == "Ada"

    def test_pipe_in_cell_text_is_escaped(self, tmp_path: Path):
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["A|B", "plain"])

        md = self._convert(tmp_path, wb)
        assert "A\\|B" in md

    def test_multiple_sheets_are_all_rendered(self, tmp_path: Path):
        import openpyxl

        wb = openpyxl.Workbook()
        wb.active.title = "First"
        wb.active.append(["x"])
        wb.create_sheet("Second").append(["y"])

        md = self._convert(tmp_path, wb)
        assert "## First" in md
        assert "## Second" in md


class TestXlrdGridBuilding:
    """Unit tests for the .xls (legacy BIFF) reader's grid construction.

    A minimal stub stands in for an xlrd Sheet rather than a real .xls file:
    xlrd can only *read* the legacy format, there is no pure-Python writer
    available in this project's dependencies (xlwt is not installed) or
    reliably available in CI (no libreoffice dependency), so this tests the
    same backfill logic MergeAwareXlsConverter._grid_from_xlrd relies on
    directly against a controlled input.
    """

    class _FakeSheet:
        def __init__(self, rows, merged_cells=()):
            self._rows = rows
            self.nrows = len(rows)
            self.ncols = len(rows[0]) if rows else 0
            self.merged_cells = merged_cells

        def cell_value(self, r, c):
            return self._rows[r][c]

    def test_backfills_a_merged_range(self):
        from converters.markitdown_xlsx import _grid_from_xlrd

        # xlrd's merged_cells format: (row_lo, row_hi, col_lo, col_hi),
        # 0-indexed with a half-open upper bound.
        sheet = self._FakeSheet(
            rows=[["Q1", ""], ["Ada", 99.0]],
            merged_cells=[(0, 1, 0, 2)],
        )
        grid = _grid_from_xlrd(sheet)
        assert grid[0] == ["Q1", "Q1"]
        assert grid[1] == ["Ada", "99"]

    def test_blank_cells_render_empty(self):
        from converters.markitdown_xlsx import _grid_from_xlrd

        sheet = self._FakeSheet(rows=[["a", ""], ["", "b"]])
        grid = _grid_from_xlrd(sheet)
        assert grid == [["a", ""], ["", "b"]]
