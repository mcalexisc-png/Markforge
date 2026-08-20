"""Programmatic fixture generation for tests.

Fixtures are generated on the fly (no binary files in the repo) so the
tests are self-contained and runnable offline.
"""

from __future__ import annotations

from pathlib import Path


def make_pdf(path: Path, *, pages: int = 3, scanned: bool = False) -> Path:
    import pymupdf

    doc = pymupdf.open()
    for index in range(1, pages + 1):
        page = doc.new_page()
        if not scanned:
            page.insert_text((72, 72), f"Chapter {index}", fontsize=24)
            page.insert_text((72, 120), f"Paragraph text on page {index}.", fontsize=12)
            page.insert_text((72, 160), "First point", fontsize=12)
            page.insert_text((72, 180), "Second point", fontsize=12)
            rect = pymupdf.Rect(72, 220, 250, 240)
            page.insert_text((72, 232), "https://example.com/page", fontsize=10)
            page.insert_link({"kind": pymupdf.LINK_URI, "from": rect, "uri": "https://example.com/page"})
    doc.set_metadata({"title": "Test Report", "author": "Markforge Tests"})
    doc.save(path)
    doc.close()
    return path


def make_scanned_pdf(path: Path, *, pages: int = 1) -> Path:
    """PDF with only images on the pages (no extractable text)."""
    from io import BytesIO

    import pymupdf
    from PIL import Image

    doc = pymupdf.open()
    for _ in range(pages):
        page = doc.new_page()
        buffer = BytesIO()
        Image.new("RGB", (400, 500), color=(240, 240, 240)).save(buffer, format="PNG")
        buffer.seek(0)
        page.insert_image(page.rect, stream=buffer.read())
    doc.save(path)
    doc.close()
    return path


def _add_hyperlink(paragraph, url: str, text: str) -> None:
    """Insert a proper external hyperlink into a docx paragraph."""
    from docx.opc.constants import RELATIONSHIP_TYPE
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    part = paragraph.part
    r_id = part.relate_to(url, RELATIONSHIP_TYPE.HYPERLINK, is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)
    run = OxmlElement("w:r")
    t = OxmlElement("w:t")
    t.text = text
    run.append(t)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def make_docx(path: Path) -> Path:
    from docx import Document

    doc = Document()
    doc.core_properties.title = "Sample Document"
    doc.core_properties.author = "Markforge Tests"
    doc.add_heading("Sample Document", level=0)
    doc.add_heading("Section One", level=1)
    doc.add_paragraph("A plain paragraph with some content.")
    p = doc.add_paragraph()
    run = p.add_run("Bold text and ")
    run.bold = True
    run2 = p.add_run("italic text")
    run2.italic = True
    p.add_run(" with a ")
    _add_hyperlink(p, "https://example.com", "link")
    doc.add_paragraph("First bullet", style="List Bullet")
    doc.add_paragraph("Second bullet", style="List Bullet")
    doc.add_paragraph("Numbered item", style="List Number")
    table = doc.add_table(rows=3, cols=2)
    table.style = "Table Grid"
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Alpha"
    table.cell(1, 1).text = "10"
    table.cell(2, 0).text = "Beta"
    table.cell(2, 1).text = "20"
    doc.add_paragraph("A quote line", style="Intense Quote")
    doc.save(path)
    return path


def make_pptx(path: Path, *, slides: int = 3) -> Path:
    from pptx import Presentation

    titles = ["Introduction", "Network Architecture", "Deployment Plan"]
    prs = Presentation()
    for index in range(1, slides + 1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = titles[index - 1]
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = f"Content paragraph for slide {index}."
        tf.add_paragraph().text = "Bullet one"
        tf.add_paragraph().text = "Bullet two"
        notes = slide.notes_slide
        notes.notes_text_frame.text = f"Speaker notes for slide {index}."
    prs.save(path)
    return path


def make_xlsx(path: Path, *, sheets: int = 2) -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Grades"
    ws.append(["Name", "Age", "Grade"])
    ws.append(["John", 20, "A"])
    ws.append(["Mary", 21, "B"])
    ws["A4"].hyperlink = "https://example.com/grade"
    ws.append(["", "22", "C"])
    ws.merge_cells("A5:B5")
    ws["A5"] = "Merged cell"
    from openpyxl.worksheet.table import Table, TableStyleInfo

    table = Table(displayName="GradesTable", ref="A1:C4")
    table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium9", showRowStripes=True)
    ws.add_table(table)

    ws2 = wb.create_sheet("Notes")
    ws2.append(["Note"])
    ws2.append(["This sheet is for notes."])
    wb.save(path)
    return path


def make_two_column_pdf(path: Path) -> Path:
    """A PDF with a full-width header and two side-by-side columns."""
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 40), "Report header spanning both columns", fontsize=11)
    for index, line in enumerate(["Left column line A", "Left column line B", "Left column line C"]):
        page.insert_text((72, 80 + index * 22), line, fontsize=11)
    for index, line in enumerate(["Right column line A", "Right column line B", "Right column line C"]):
        page.insert_text((340, 80 + index * 22), line, fontsize=11)
    doc.save(path)
    doc.close()
    return path


def make_deck_pdf(path: Path, *, missions: int = 3) -> Path:
    """A deck-style PDF: large-font mission titles, two-column pages,
    one duplicated page (animation build) and one image-only page."""
    from io import BytesIO

    import pymupdf
    from PIL import Image

    doc = pymupdf.open()
    for mission in range(1, missions + 1):
        page = doc.new_page()
        page.insert_text((72, 60), f"Mission {mission}", fontsize=32)
        for column, prefix in ((72, "Left"), (340, "Right")):
            for line_index in range(3):
                page.insert_text(
                    (column, 110 + line_index * 24),
                    f"{prefix} column line {line_index + 1} of mission {mission}",
                    fontsize=12,
                )
    doc.insert_pdf(
        pymupdf.open(stream=doc.tobytes(), filetype="pdf"),
        from_page=missions - 1,
        to_page=missions - 1,
    )

    page = doc.new_page()
    buffer = BytesIO()
    Image.new("RGB", (400, 500), color=(240, 240, 240)).save(buffer, format="PNG")
    buffer.seek(0)
    page.insert_image(page.rect, stream=buffer.read())

    doc.save(path)
    doc.close()
    return path